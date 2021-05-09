from pptx import Presentation
import os
import shutil

src = r"E:\火车头\总文件"
src1 = r"E:\火车头\short_file"
dir_list = os.listdir(src)
print(dir_list)
for dir in dir_list:
    try:
        prs = Presentation('{}/{}'.format(src, dir))
        slides = prs.slides
        number_pages = len(slides)
        print(number_pages)
        # 读取ppt
        # if number_pages <= 8:
        #     shutil.move(r'E:\火车头\ppt_pptx\{}'.format(dir), r'E:\火车头\short_file\{}'.format(dir))
        #     # continue

        # 删除最后一页
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]
        #
        # rId = prs.slides._sldIdLst[-1].rId
        # prs.part.drop_rel(rId)
        # del prs.slides._sldIdLst[-1]
        #
        # rId = prs.slides._sldIdLst[-1].rId
        # prs.part.drop_rel(rId)
        # del prs.slides._sldIdLst[-1]
        #
        # rId = prs.slides._sldIdLst[-1].rId
        # prs.part.drop_rel(rId)
        # del prs.slides._sldIdLst[-1]
        #
        # rId = prs.slides._sldIdLst[-1].rId
        # prs.part.drop_rel(rId)
        # del prs.slides._sldIdLst[-1]
        # 保存新的ppt
        prs.save('{}\{}'.format(src1, dir))
    except:
        print("出错了")
        try:
            shutil.move(r'E:\火车头\总文件\{}'.format(dir), r'E:\火车头\err\{}'.format(dir))
        except:
            print("复制文件异常")
            print(dir)

# prs = Presentation('./out_ppt3/pptx')
# # 查看一共几页
# slides = prs.slides
# number_pages = len(slides)
# print(number_pages)
# 删除最后一页
# rId = prs.slides._sldIdLst[-1].rId
# prs.part.drop_rel(rId)
# del prs.slides._sldIdLst[-1]
# # 保存新的ppt
# prs.save('./new.pptx')

# import os
# import shutil
#
# dir_list = os.listdir("./out_ppt2")
# print(dir_list)
#
# for dir_name in dir_list:
#     srcFile = './out_ppt2/{}'.format(dir_name)
#     # print(dir_name[-4:])
#     if dir_name[-3:] == "ppt":
#         dstFile = './out_ppt3/ppt/{}'.format(dir_name)
#         shutil.move(srcFile, dstFile)
#     elif dir_name[-4:] == "pptx":
#         dstFile = './out_ppt3/pptx/{}'.format(dir_name)
#         shutil.move(srcFile, dstFile)
#     else:
#         pass


# new_name = str.replace(dir_name, "-www.1ppt.com.ppt", "")
# new_name = str.replace(new_name, "www.1ppt.com", "")
# new_name = str.replace(new_name, "（", "")
# new_name = str.replace(new_name, "）", "")
# new_name = str.replace(new_name, "【", "")
# new_name = str.replace(new_name, "】", "")
# new_name = str.replace(new_name, "[", "")
# new_name = str.replace(new_name, "]", "")
# new_name = str.replace(new_name, "(", "")
# new_name = str.replace(new_name, ")", "")
# print(new_name)
#
# dstFile = './out_ppt2/{}'.format(new_name)
# try:
#     os.rename(srcFile, dstFile)
# except Exception as e:
#     print(e)
#     print('rename file fail\r\n')

# try:
#     os.rename(srcFile,dstFile)
# except Exception as e:
#     print(e)
#     print('rename file fail\r\n')
# else:
#     print('rename file success\r\n')
